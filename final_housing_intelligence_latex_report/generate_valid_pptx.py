from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "housing_intelligence_25_slides.pptx"

NAVY = RGBColor(23, 50, 77)
TEAL = RGBColor(8, 127, 140)
ORANGE = RGBColor(195, 95, 29)
MUTED = RGBColor(83, 101, 122)
BG = RGBColor(247, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(216, 225, 234)


def add_background(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def add_top_bar(slide, title: str):
    add_background(slide)
    tag = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(3.3), Inches(0.25))
    p = tag.text_frame.paragraphs[0]
    p.text = "US HOUSING INTELLIGENCE"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.48), Inches(12.0), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.15), Inches(12.4), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, size: int = 18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED
        p.space_after = Pt(8)


def add_picture(slide, image_path: Path, left: float, top: float, width: float | None = None, height: float | None = None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    pic = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)
    pic.line.color.rgb = LINE
    pic.line.width = Pt(0.75)
    return pic


def title_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    tag = slide.shapes.add_textbox(Inches(0.65), Inches(0.62), Inches(7.0), Inches(0.35))
    p = tag.text_frame.paragraphs[0]
    p.text = "US HOUSING INTELLIGENCE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(190, 239, 239)

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(11.8), Inches(1.25))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_bullets(slide, bullets, 0.82, 3.0, 9.8, 2.3, 20)
    footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.55), Inches(9.5), Inches(0.35))
    p = footer.text_frame.paragraphs[0]
    p.text = "Data Mining and Machine Learning | Final Defense"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(221, 232, 240)


def bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, title)
    add_bullets(slide, bullets, 0.8, 1.65, 11.6, 5.2, 19)


def image_full_slide(prs: Presentation, title: str, image: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, title)
    add_picture(slide, image, 0.85, 1.45, width=11.6)


def image_left_slide(prs: Presentation, title: str, image: Path, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, title)
    add_picture(slide, image, 0.65, 1.55, width=7.0)
    add_bullets(slide, bullets, 8.0, 1.7, 4.7, 4.9, 16)


def metrics_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, title)
    metrics = [("242", "Monthly observations", TEAL), ("27", "Dataset columns", ORANGE), ("2004-2024", "Historical period", NAVY)]
    for idx, (value, label, color) in enumerate(metrics):
        left = 0.9 + idx * 4.1
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.15), Inches(1.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE

        lab = slide.shapes.add_textbox(Inches(left), Inches(3.16), Inches(3.15), Inches(0.4))
        p = lab.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
    add_bullets(slide, bullets, 1.0, 4.05, 11.1, 2.2, 17)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    img = {
        "home": ROOT / "app_screenshots" / "app_home.png",
        "modeling": ROOT / "app_screenshots" / "app_modeling.png",
        "evaluation_app": ROOT / "app_screenshots" / "app_evaluation.png",
        "olap_app": ROOT / "app_screenshots" / "app_olap.png",
        "paper_app": ROOT / "app_screenshots" / "app_paper_review.png",
        "governance": ROOT / "app_screenshots" / "app_governance.png",
        "pipeline": ROOT / "figures" / "project_pipeline.png",
        "corr": ROOT / "figures" / "correlation_results.png",
        "eval": ROOT / "figures" / "model_evaluation_results.png",
        "diag": ROOT / "figures" / "fit_diagnostics_results.png",
        "pred": ROOT / "figures" / "actual_vs_predicted_results.png",
        "olap": ROOT / "figures" / "olap_period_results.png",
    }

    title_slide(prs, "US Housing Intelligence Platform", [
        "Final report and Streamlit app defense",
        "Data Mining and Machine Learning project",
        "Ayoub Naouech and Rayen Belhadj",
    ])
    bullet_slide(prs, "Presentation Roadmap", [
        "Context, problem, and objectives",
        "Dataset, preparation, and CRISP-DM methodology",
        "Streamlit app screenshots and modules",
        "Evaluation, diagnostics, OLAP, and paper-review evidence",
        "Limitations, future work, and conclusion",
    ])
    bullet_slide(prs, "Project Context", [
        "The U.S. housing market connects household wealth, credit conditions, inflation, construction, and employment.",
        "Housing prices react to demand, supply, income, mortgage rates, expectations, and timing.",
        "The project converts raw housing and macroeconomic data into an interactive intelligence platform.",
    ])
    bullet_slide(prs, "Problem Statement", [
        "Main question: how can data mining and machine learning analyze and predict U.S. housing-market trends?",
        "Analytical challenge: housing variables interact in non-linear and time-dependent ways.",
        "Communication challenge: model results must be understandable for non-technical users.",
    ])
    bullet_slide(prs, "Project Objectives", [
        "Analyze historical evolution of the Home Price Index.",
        "Identify variables associated with U.S. housing-price behavior.",
        "Train and compare machine learning models using reproducible metrics.",
        "Build a Streamlit app for visualization, modeling, OLAP, export, and interpretation.",
    ])
    metrics_slide(prs, "Dataset Overview", [
        "Target variable: Home_Price_Index.",
        "Main variables: mortgage rate, interest rate, unemployment, inflation, permits, income, population, and sentiment.",
        "The dataset supports both prediction and interpretation.",
    ])
    bullet_slide(prs, "Data Preparation", [
        "Parse and sort the date column chronologically.",
        "Convert numeric columns into model-ready values.",
        "Check missing values and duplicate rows before modeling.",
        "Create administration labels, lag features, rolling features, ratios, and percentage changes.",
    ])
    bullet_slide(prs, "CRISP-DM Methodology", [
        "Business understanding defines the housing problem.",
        "Data understanding checks quality and temporal structure.",
        "Data preparation builds usable analytical features.",
        "Modeling and evaluation compare algorithms.",
        "Deployment is delivered through Streamlit.",
    ])
    image_left_slide(prs, "System Architecture", img["pipeline"], [
        "Data layer: loading, cleaning, validation, and features.",
        "Analysis layer: EDA, correlations, OLAP, and period comparison.",
        "Modeling layer: regression, classification, forecasting, and diagnostics.",
        "Interpretation layer: paper review, conclusion, and export.",
    ])
    image_full_slide(prs, "Streamlit App: Home And Workflow", img["home"])
    image_full_slide(prs, "Streamlit App: ML Lab", img["modeling"])
    image_full_slide(prs, "Streamlit App: Evaluation", img["evaluation_app"])
    image_full_slide(prs, "Streamlit App: OLAP And Export", img["olap_app"])
    image_left_slide(prs, "Streamlit App: Paper Review", img["paper_app"], [
        "Compares dashboard evidence with housing theory.",
        "Explains rates, supply, demand, timing, and affordability together.",
        "Makes the project academically stronger than a purely technical dashboard.",
    ])
    image_full_slide(prs, "Streamlit App: Governance", img["governance"])
    image_left_slide(prs, "Correlation Results", img["corr"], [
        "Strong relationships appear around lagged HPI, smoothed HPI, and Median_Income.",
        "Correlation supports interpretation but does not prove causality.",
        "Housing prices are persistent and linked to broader economic context.",
    ])
    image_left_slide(prs, "Model Evaluation Results", img["eval"], [
        "Compared models include Linear Regression, Ridge, Random Forest, Gradient Boosting, Decision Tree, KNN, and SVR.",
        "Best generated metrics: R2 = 0.999, RMSE = 0.797, MAE = 0.661.",
    ])
    bullet_slide(prs, "Regression Metrics Table", [
        "Linear Regression: MAE 0.661, RMSE 0.797, R2 0.999, CV R2 0.994.",
        "Ridge Regression: MAE 3.715, RMSE 4.162, R2 0.984.",
        "Tree-based models perform weaker on the final chronological test period.",
        "Evaluation must be time-aware, not based only on model popularity.",
    ])
    image_left_slide(prs, "Fit Diagnostics", img["diag"], [
        "Diagnostics compare train and test performance.",
        "Large train-test gaps indicate weak generalization or overfitting risk.",
        "Tree models struggle when the final period moves beyond the training range.",
    ])
    image_left_slide(prs, "Actual Versus Predicted", img["pred"], [
        "The best model follows the chronological test period closely.",
        "Lagged and smoothed price features explain the strong predictive result.",
        "The visual is easier to defend than metrics alone.",
    ])
    image_left_slide(prs, "OLAP Period Results", img["olap"], [
        "Biden period has the highest average Home Price Index.",
        "Biden period also has the highest average mortgage-rate pressure.",
        "Prices and affordability pressure can rise together when supply and timing matter.",
    ])
    bullet_slide(prs, "Paper Review Message", [
        "Housing prices are shaped by demand, financing conditions, supply, construction, income, and timing.",
        "Higher mortgage rates reduce affordability, but prices may react slowly when supply is constrained.",
        "The strongest conclusion is multi-factor interpretation rather than one-variable explanation.",
    ])
    bullet_slide(prs, "Business Impact", [
        "Analytical impact: users understand which indicators move with housing prices.",
        "Decision-support impact: users compare models, periods, OLAP segments, and scenarios.",
        "Educational impact: the project demonstrates EDA, ML, forecasting, OLAP, paper review, and governance.",
    ])
    bullet_slide(prs, "Limitations And Future Work", [
        "National-level data hides regional differences between states and cities.",
        "Forecasts are educational and should not be treated as financial advice.",
        "Future work: regional datasets, live APIs, SHAP explainability, better chatbot grounding, monitoring, and authentication.",
    ])
    bullet_slide(prs, "Final Conclusion", [
        "The report, poster, and app now tell one consistent story.",
        "U.S. housing prices require theory, correlations, model evaluation, diagnostics, OLAP, and market context.",
        "The Streamlit dashboard is a housing intelligence platform, not only a collection of charts.",
    ])

    prs.save(OUT)
    print(f"Created {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
